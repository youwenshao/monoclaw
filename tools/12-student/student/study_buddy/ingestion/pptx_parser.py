"""PowerPoint document parser using python-pptx."""

from __future__ import annotations

from pptx import Presentation


def parse_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    slides: list[dict] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        texts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.shape_type == 13 or (hasattr(shape, "is_placeholder") and shape.placeholder_format and shape.placeholder_format.idx == 0):
                    title = shape.text_frame.text.strip()
                texts.append(shape.text_frame.text.strip())

        if not title and texts:
            title = texts[0][:80]

        combined = "\n".join(texts).strip()
        if combined:
            slides.append({
                "slide_number": idx,
                "text": combined,
                "title": title,
            })

    return slides
